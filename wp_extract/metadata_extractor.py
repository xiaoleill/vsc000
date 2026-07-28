"""
Metadata Extractor Module
Extracts Document No., Revision, Author, Reviewer, Approver, Status, Date
from WP documents in various formats.

Supported formats:
- Excel (.xlsx, .xlsm): Reads "Title" or "Cover" sheet, B/C column key-value pairs
- Word (.docx): Reads docProps and first table
- PowerPoint (.pptx): Reads docProps and first slide text
- PDF: Extracts first page text and parses metadata patterns
- Markdown (.md): Parses YAML frontmatter or metadata patterns
"""

import os
import re
from dataclasses import dataclass, field
from datetime import datetime
from typing import Optional

# Metadata field labels to search for (case insensitive, in priority order)
FIELD_PATTERNS = {
    'document_no': [
        r'\bdocument\s*no\.?\s*[:：]?\s*$',
        r'\bdocument\s*number\s*[:：]?\s*$',
        r'\bdoc\s*no\.?\s*[:：]?\s*$',
        r'\bdoc\s*id\s*[:：]?\s*$',
        r'编号\s*[:：]?\s*$',
        r'文檔編號\s*[:：]?\s*$',
    ],
    'revision': [
        r'\brevision\s*[:：]?\s*$',
        r'\bversion\s*[:：]?\s*$',
        r'\bver\.?\s*[:：]?\s*$',
        r'版本\s*[:：]?\s*$',
    ],
    'author': [
        r'\bauthor\s*[:：]?\s*$',
        r'\bprepared\s*by\s*[:：]?\s*$',
        r'\bcreated\s*by\s*[:：]?\s*$',
        r'作者\s*[:：]?\s*$',
        r'编制\s*[:：]?\s*$',
        r'编\s*制\s*[:：]?\s*$',
        r'編制\s*[:：]?\s*$',
    ],
    'reviewer': [
        r'\breviewer\s*[:：]?\s*$',
        r'\breviewed\s*by\s*[:：]?\s*$',
        r'\bchecked\s*by\s*[:：]?\s*$',
        r'审核\s*[:：]?\s*$',
        r'审\s*核\s*[:：]?\s*$',
        r'審核\s*[:：]?\s*$',
    ],
    'approver': [
        r'\bapprover\s*[:：]?\s*$',
        r'\bapproved\s*by\s*[:：]?\s*$',
        r'批准\s*[:：]?\s*$',
        r'批\s*准\s*[:：]?\s*$',
        r'核准\s*[:：]?\s*$',
    ],
    'status': [
        r'\bstatus\s*[:：]?\s*$',
        r'\bstate\s*[:：]?\s*$',
        r'状态\s*[:：]?\s*$',
        r'狀態\s*[:：]?\s*$',
    ],
    'date': [
        r'\bdate\s*[:：]?\s*$',
        r'\blast\s*modified\s*[:：]?\s*$',
        r'日期\s*[:：]?\s*$',
    ],
}

NOT_FOUND = "NA"


@dataclass
class DocumentMetadata:
    """Extracted metadata from a WP document."""
    file_path: str
    document_no: str = NOT_FOUND
    revision: str = NOT_FOUND
    author: str = NOT_FOUND
    reviewer: str = NOT_FOUND
    approver: str = NOT_FOUND
    status: str = NOT_FOUND
    date: str = NOT_FOUND

    def to_dict(self) -> dict:
        return {
            'file_path': self.file_path,
            'document_no': self.document_no,
            'revision': self.revision,
            'author': self.author,
            'reviewer': self.reviewer,
            'approver': self.approver,
            'status': self.status,
            'date': self.date,
        }


def _match_field_label(text: str) -> Optional[str]:
    """
    Check if text matches any known field label pattern.
    Returns the field name (key in FIELD_PATTERNS) or None.
    """
    text_clean = text.strip().lower()
    for field_name, patterns in FIELD_PATTERNS.items():
        for pattern in patterns:
            if re.search(pattern, text_clean):  # search, not match — label may have prefix
                return field_name
    return None


def _parse_date_value(text: str) -> str:
    """Normalize date string to YYYY-MM-DD format."""
    text = str(text).strip() if text is not None else ''
    if not text:
        return NOT_FOUND

    # If it's already a datetime, format it
    if hasattr(text, 'strftime'):
        return text.strftime('%Y-%m-%d')

    # Strip time portion if present: "2024-06-15 00:00:00" -> "2024-06-15"
    if ' ' in text:
        text = text.split(' ')[0]

    # Try various date formats, normalize to YYYY-MM-DD
    import re as _re
    # YYYY/M/D or YYYY-MM-DD or YYYY.M.D
    m = _re.match(r'(\d{4})[/\-\.](\d{1,2})[/\-\.](\d{1,2})', text)
    if m:
        return f'{m.group(1)}-{int(m.group(2)):02d}-{int(m.group(3)):02d}'

    return text


def _normalize_value(text: str) -> str:
    """Normalize a metadata value string."""
    if text is None:
        return NOT_FOUND
    if isinstance(text, datetime):
        return text.strftime('%Y-%m-%d')
    text = str(text).strip()
    if not text:
        return NOT_FOUND
    # Remove excessive whitespace
    text = ' '.join(text.split())
    return text


def _extract_from_key_value_pairs(pairs: list[tuple[str, str]]) -> dict:
    """
    Given a list of (label, value) pairs, extract metadata fields.
    """
    result = {
        'document_no': NOT_FOUND,
        'revision': NOT_FOUND,
        'author': NOT_FOUND,
        'reviewer': NOT_FOUND,
        'approver': NOT_FOUND,
        'status': NOT_FOUND,
        'date': NOT_FOUND,
    }

    for label, value in pairs:
        field = _match_field_label(label)
        if field:
            if field == 'date':
                result[field] = _parse_date_value(value)
            else:
                result[field] = _normalize_value(value)

    return result


# ============ Excel Extraction ============

def extract_excel_metadata(filepath: str) -> DocumentMetadata:
    """
    Extract metadata from Excel file's Title or Cover sheet.
    Reads B/C column key-value pairs from the first matching sheet.
    """
    try:
        import openpyxl
    except ImportError:
        return DocumentMetadata(file_path=filepath)

    meta = DocumentMetadata(file_path=filepath)

    wb = None
    # Try normal mode first (handles formulas, merged cells)
    try:
        wb = openpyxl.load_workbook(filepath, data_only=True)
    except Exception:
        pass

    # Fallback: read_only mode (less strict XML validation)
    if wb is None:
        try:
            wb = openpyxl.load_workbook(filepath, read_only=True, data_only=True)
        except Exception:
            return meta

    # Priority order for sheet names
    target_sheets = []
    for sn in wb.sheetnames:
        sn_lower = sn.lower().strip()
        if sn_lower in ('title', 'cover', '封面', '文档信息', 'document info'):
            target_sheets.append(sn)

    # If no standard sheet found, try the first sheet
    if not target_sheets and wb.sheetnames:
        target_sheets.append(wb.sheetnames[0])

    for sn in target_sheets:
        ws = wb[sn]
        pairs = []

        # Read all rows looking for B/C column key-value pairs
        for row in ws.iter_rows(min_row=1, max_row=min(50, ws.max_row),
                                min_col=2, max_col=3, values_only=True):
            if len(row) >= 2:
                label = str(row[0]).strip() if row[0] is not None else ""
                value = row[1]
                if label:
                    value_str = str(value) if value is not None else ""
                    pairs.append((label, value_str))

        fields = _extract_from_key_value_pairs(pairs)

        # Accumulate results — fill gaps without overwriting
        for k, v in fields.items():
            if getattr(meta, k) == NOT_FOUND and v != NOT_FOUND:
                setattr(meta, k, v)

    # Also check 'Change History' sheet for Revision/Author/Date
    for sn in wb.sheetnames:
        sn_lower = sn.strip().lower()
        if 'change history' not in sn_lower and 'revision history' not in sn_lower:
            continue
        ws = wb[sn]
        # Find the header row (usually row 2: Revision | Date | Author | ...)
        header_row = None
        rev_col = date_col = author_col = None
        for r in range(1, min(5, ws.max_row + 1)):
            for c in range(1, min(6, ws.max_column + 1)):
                v = str(ws.cell(row=r, column=c).value or '').strip().lower()
                if v == 'revision':
                    rev_col = c
                    header_row = r
                elif v == 'date':
                    date_col = c
                elif v == 'author':
                    author_col = c
        if rev_col and header_row:
            # Read all data rows, take the LAST one with a non-empty revision
            last_rev = last_date = last_author = None
            for r in range(header_row + 1, ws.max_row + 1):
                rev_val = ws.cell(row=r, column=rev_col).value
                if rev_val is not None and str(rev_val).strip():
                    last_rev = str(rev_val).strip()
                    if date_col:
                        dv = ws.cell(row=r, column=date_col).value
                        last_date = str(dv).strip() if dv is not None else None
                    if author_col:
                        av = ws.cell(row=r, column=author_col).value
                        last_author = str(av).strip() if av is not None else None
            if meta.revision == NOT_FOUND and last_rev:
                meta.revision = last_rev
            if meta.author == NOT_FOUND and last_author:
                meta.author = last_author
            if meta.date == NOT_FOUND and last_date:
                meta.date = last_date
            break  # Only process first Change History sheet

    wb.close()
    return meta


# ============ Word Extraction ============

def _cell_text_including_sdt(cell) -> str:
    """
    Get full text from a table cell, INCLUDING text inside SDT content controls.
    python-docx's cell.text skips <w:sdt> content, so we parse the XML directly.
    """
    # First try python-docx default
    base = cell.text.strip()
    if base:
        return base

    # Fallback: extract text from all <w:t> elements including those inside SDTs
    NS = '{http://schemas.openxmlformats.org/wordprocessingml/2006/main}'
    texts = []
    for t_elem in cell._tc.iter(f'{NS}t'):
        if t_elem.text:
            texts.append(t_elem.text)
    return ''.join(texts).strip()


def _extract_sdt_pairs(doc) -> list[tuple[str, str]]:
    """
    Extract label-value pairs from Structured Document Tags (SDT) in a docx.
    SDTs are Word content controls often used for cover page metadata.
    Parses <w:sdt> elements in document body, extracting <w:t> text runs
    and splitting concatenated label+value sequences.
    """
    pairs = []
    # XML namespace for Word
    NS = '{http://schemas.openxmlformats.org/wordprocessingml/2006/main}'

    for child in doc.element.body:
        tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
        if tag != 'sdt':
            continue

        # Collect all text from <w:t> elements within this SDT
        texts = []
        for t_elem in child.iter(f'{NS}t'):
            if t_elem.text:
                texts.append(t_elem.text)

        full_text = ''.join(texts)

        # Use known label patterns to split the concatenated text
        # Order matters: match longer labels first to avoid partial matches
        label_patterns = [
            (r'(Document\s*No\.?)', 'document_no'),
            (r'(Revision)', 'revision'),
            (r'(Confidential\s*level)', None),       # skip, not a target field
            (r'(Author)', 'author'),
            (r'(Reviewer)', 'reviewer'),
            (r'(Approver)', 'approver'),
            (r'(Status)', 'status'),
            (r'(Date)', 'date'),
            (r'(Title\s*[:：]?)', None),              # skip document title
            (r'(Subject)', None),
            (r'(编号)', 'document_no'),
            (r'(版本)', 'revision'),
            (r'(作者)', 'author'),
            (r'(审核)', 'reviewer'),
            (r'(批准)', 'approver'),
            (r'(状态)', 'status'),
            (r'(日期)', 'date'),
            (r'(编制)', 'author'),
            (r'(审\s*核)', 'reviewer'),
            (r'(批\s*准)', 'approver'),
        ]

        # Build regex that matches any known label at word boundaries
        combined = '|'.join(f'(?P<f{i}>{p})' for i, (p, _) in enumerate(label_patterns))
        pattern = re.compile(combined)

        # Find all label positions
        matches = list(pattern.finditer(full_text))
        if len(matches) < 2:
            continue  # need at least 2 labels to extract meaningful pairs

        for i, m in enumerate(matches):
            # Determine which field matched
            matched_field = None
            for j, (_, field_name) in enumerate(label_patterns):
                if m.group(f'f{j}'):
                    matched_field = field_name
                    break

            if matched_field is None:
                continue  # skip non-target labels (e.g., Confidential level)

            label_text = m.group(0)
            # Value is text between this label end and next label start
            value_start = m.end()
            value_end = matches[i + 1].start() if i + 1 < len(matches) else len(full_text)
            value = full_text[value_start:value_end].strip()
            if value:
                pairs.append((label_text, value))

            # Stop after Date — last cover page field; rest is revision history
            if matched_field == 'date':
                return pairs

        break  # Only process the first SDT (cover page)

    return pairs


def extract_docx_metadata(filepath: str) -> DocumentMetadata:
    """
    Extract metadata from Word document.
    Search order (front-to-back priority):
      1. SDT elements (cover page content controls)
      2. First 2 tables in body
      3. First 30 paragraphs (Label: Value pattern)
      4. docProps (lowest priority fallback)
    """
    try:
        from docx import Document
    except ImportError:
        return DocumentMetadata(file_path=filepath)

    meta = DocumentMetadata(file_path=filepath)

    try:
        doc = Document(filepath)
    except Exception:
        return meta

    fields = {
        'document_no': NOT_FOUND,
        'revision': NOT_FOUND,
        'author': NOT_FOUND,
        'reviewer': NOT_FOUND,
        'approver': NOT_FOUND,
        'status': NOT_FOUND,
        'date': NOT_FOUND,
    }

    def _apply(pairs):
        """Apply pairs to fields dict, only filling gaps (no overwrite)."""
        extracted = _extract_from_key_value_pairs(pairs)
        for k in fields:
            if fields[k] == NOT_FOUND and extracted[k] != NOT_FOUND:
                fields[k] = extracted[k]

    # Priority 1: SDT elements (cover page content controls — earliest content)
    sdt_pairs = _extract_sdt_pairs(doc)
    if sdt_pairs:
        _apply(sdt_pairs)

    # Priority 2: First 2 tables in body
    for table in doc.tables[:2]:
        pairs = []
        for row in table.rows:
            cells = row.cells
            if len(cells) >= 2:
                label = _cell_text_including_sdt(cells[0])
                value = _cell_text_including_sdt(cells[1])
                if label:
                    pairs.append((label, value))
        _apply(pairs)

    # Priority 3: First 30 paragraphs (Label: Value pattern)
    para_pairs = []
    for para in doc.paragraphs[:30]:
        text = para.text.strip()
        if ':' in text or '：' in text:
            parts = re.split(r'[:：]', text, maxsplit=1)
            if len(parts) == 2:
                label = parts[0].strip()
                value = parts[1].strip()
                if label and value and len(label) < 50:
                    para_pairs.append((label, value))
    _apply(para_pairs)

    # Priority 4: docProps (fallback only)
    try:
        cp = doc.core_properties
        if fields['author'] == NOT_FOUND and cp.author:
            fields['author'] = cp.author
        if fields['reviewer'] == NOT_FOUND and cp.last_modified_by:
            fields['reviewer'] = cp.last_modified_by
        if fields['date'] == NOT_FOUND and cp.modified:
            fields['date'] = cp.modified.strftime('%Y-%m-%d')
    except Exception:
        pass

    meta.document_no = fields['document_no']
    meta.revision = fields['revision']
    meta.author = fields['author']
    meta.reviewer = fields['reviewer']
    meta.approver = fields['approver']
    meta.status = fields['status']
    meta.date = fields['date']

    return meta


# ============ PowerPoint Extraction ============

def extract_pptx_metadata(filepath: str) -> DocumentMetadata:
    """
    Extract metadata from PowerPoint document.
    Search order (front-to-back priority):
      1. First slide shapes (tables then text frames)
      2. docProps (fallback only)
    """
    try:
        from pptx import Presentation
    except ImportError:
        return DocumentMetadata(file_path=filepath)

    meta = DocumentMetadata(file_path=filepath)

    try:
        prs = Presentation(filepath)
    except Exception:
        return meta

    fields = {
        'document_no': NOT_FOUND,
        'revision': NOT_FOUND,
        'author': NOT_FOUND,
        'reviewer': NOT_FOUND,
        'approver': NOT_FOUND,
        'status': NOT_FOUND,
        'date': NOT_FOUND,
    }

    def _apply(pairs):
        extracted = _extract_from_key_value_pairs(pairs)
        for k in fields:
            if fields[k] == NOT_FOUND and extracted[k] != NOT_FOUND:
                fields[k] = extracted[k]

    # Priority 1: First slide — tables first, then text frames
    if prs.slides:
        slide = prs.slides[0]
        for shape in slide.shapes:
            if shape.has_table:
                pairs = []
                table = shape.table
                for row in table.rows:
                    cells = row.cells
                    if len(cells) >= 2:
                        label = cells[0].text.strip()
                        value = cells[1].text.strip()
                        if label:
                            pairs.append((label, value))
                _apply(pairs)
            elif shape.has_text_frame:
                pairs = []
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if ':' in text or '：' in text:
                        parts = re.split(r'[:：]', text, maxsplit=1)
                        if len(parts) == 2:
                            label = parts[0].strip()
                            value = parts[1].strip()
                            if label and value and len(label) < 50:
                                pairs.append((label, value))
                _apply(pairs)

    # Priority 2: docProps (fallback only)
    try:
        cp = prs.core_properties
        if fields['author'] == NOT_FOUND and cp.author:
            fields['author'] = cp.author
        if fields['reviewer'] == NOT_FOUND and cp.last_modified_by:
            fields['reviewer'] = cp.last_modified_by
        if fields['date'] == NOT_FOUND and cp.modified:
            fields['date'] = cp.modified.strftime('%Y-%m-%d')
    except Exception:
        pass

    meta.document_no = fields['document_no']
    meta.revision = fields['revision']
    meta.author = fields['author']
    meta.reviewer = fields['reviewer']
    meta.approver = fields['approver']
    meta.status = fields['status']
    meta.date = fields['date']

    return meta


# ============ PDF Extraction ============

def extract_pdf_metadata(filepath: str) -> DocumentMetadata:
    """
    Extract metadata from PDF document.
    Uses pdfplumber for text extraction, falls back to PyMuPDF.
    """
    meta = DocumentMetadata(file_path=filepath)

    text = ""

    # Try pdfplumber first
    try:
        import pdfplumber
        with pdfplumber.open(filepath) as pdf:
            if pdf.pages:
                # Extract first page text
                page = pdf.pages[0]
                t = page.extract_text()
                if t:
                    text = t
                # Also try table extraction
                tables = page.extract_tables()
                if tables:
                    for table in tables:
                        for row in table:
                            if row and len(row) >= 2:
                                label = str(row[0] or '').strip()
                                value = str(row[1] or '').strip()
                                if label:
                                    text += f'\n{label}: {value}'
    except ImportError:
        pass
    except Exception:
        pass

    # Fallback to PyMuPDF
    if not text:
        try:
            import fitz  # PyMuPDF
            doc = fitz.open(filepath)
            if doc.page_count > 0:
                page = doc[0]
                text = page.get_text()
            doc.close()
        except ImportError:
            pass
        except Exception:
            pass

    if not text:
        return meta

    # Parse text for metadata patterns
    pairs = []
    lines = text.split('\n')
    for line in lines:
        line = line.strip()
        if not line:
            continue
        if ':' in line or '：' in line:
            parts = re.split(r'[:：]', line, maxsplit=1)
            if len(parts) == 2:
                label = parts[0].strip()
                value = parts[1].strip()
                if label and value and len(label) < 50:
                    pairs.append((label, value))

    fields = _extract_from_key_value_pairs(pairs)

    meta.document_no = fields['document_no']
    meta.revision = fields['revision']
    meta.author = fields['author']
    meta.reviewer = fields['reviewer']
    meta.approver = fields['approver']
    meta.status = fields['status']
    meta.date = fields['date']

    return meta


# ============ Markdown Extraction ============

def extract_md_metadata(filepath: str) -> DocumentMetadata:
    """
    Extract metadata from Markdown file.
    Parses YAML frontmatter and metadata patterns in text.
    """
    meta = DocumentMetadata(file_path=filepath)

    try:
        with open(filepath, 'r', encoding='utf-8') as f:
            content = f.read()
    except Exception:
        try:
            with open(filepath, 'r', encoding='gbk') as f:
                content = f.read()
        except Exception:
            return meta

    pairs = []

    # Check for YAML frontmatter
    fm_match = re.match(r'^---\s*\n(.*?)\n---', content, re.DOTALL)
    if fm_match:
        frontmatter = fm_match.group(1)
        for line in frontmatter.split('\n'):
            if ':' in line:
                parts = line.split(':', 1)
                if len(parts) == 2:
                    label = parts[0].strip()
                    value = parts[1].strip().strip('"').strip("'")
                    if label and value:
                        pairs.append((label, value))

    # Also search content for metadata patterns
    for line in content.split('\n')[:50]:
        line = line.strip()
        if not line or line.startswith('#'):
            continue
        # Pattern: **Label:** value or **Label** value
        bold_match = re.match(r'\*\*(.+?)\*\*\s*[:：]?\s*(.+)', line)
        if bold_match:
            label = bold_match.group(1).strip().rstrip(':').rstrip('：')
            value = bold_match.group(2).strip()
            if label and value:
                pairs.append((label, value))

    fields = _extract_from_key_value_pairs(pairs)

    meta.document_no = fields['document_no']
    meta.revision = fields['revision']
    meta.author = fields['author']
    meta.reviewer = fields['reviewer']
    meta.approver = fields['approver']
    meta.status = fields['status']
    meta.date = fields['date']

    return meta


# ============ Dispatcher ============

EXTRACTORS = {
    '.xlsx': extract_excel_metadata,
    '.xlsm': extract_excel_metadata,
    '.docx': extract_docx_metadata,
    '.pptx': extract_pptx_metadata,
    '.pdf': extract_pdf_metadata,
    '.md': extract_md_metadata,
}


def extract_metadata(filepath: str) -> DocumentMetadata:
    """
    Extract metadata from any supported file format.
    Dispatches to the appropriate format-specific extractor.
    """
    if not os.path.exists(filepath):
        return DocumentMetadata(file_path=filepath)

    ext = os.path.splitext(filepath)[1].lower()
    extractor = EXTRACTORS.get(ext)

    if extractor is None:
        return DocumentMetadata(file_path=filepath)

    return extractor(filepath)
